"""Assemble the five-slide comparison PPTX."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

from bedrock.analysis.electricity.current.diagnostics.deck.data import (
    ImplBundle,
    fill_mixed_c_col,
)
from bedrock.analysis.electricity.current.diagnostics.deck.histograms import (
    write_hist_png,
)
from bedrock.analysis.electricity.current.diagnostics.deck.pairs import Pair
from bedrock.analysis.electricity.current.diagnostics.deck.tables import (
    TableGrid,
    class_mwh_grids,
    ef_grids,
    qx_grids,
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
# Petrol header from the original vs EIA-anchored deck (≈ RGB 0, 91, 127).
HEADER_BG = RGBColor(0x00, 0x5B, 0x7F)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ALT_ROW = RGBColor(0xEB, 0xEB, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_FG = RGBColor(0x22, 0x22, 0x22)
SAME_FG = RGBColor(0x2E, 0x7D, 0x32)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
TABLE_FONT_PT = 12
TABLE_WIDTH_IN = 6.15
TABLE_LEFTS = (0.40, 6.78)
TABLE_TOP_IN = 3.20
HEADER_ROW_IN = 0.34
DATA_ROW_IN = 0.30
BORDER_WHITE = 'FFFFFF'
BORDER_EMU = '12700'


def _set_run_font(
    run: Any, *, size_pt: float, bold: bool = False, color: RGBColor | None = None
) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = 'Calibri'
    if color is not None:
        run.font.color.rgb = color


def _set_cell_borders(cell: Any) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    for edge in ('lnL', 'lnR', 'lnT', 'lnB'):
        for child in list(tc_pr):
            if child.tag.endswith(edge):
                tc_pr.remove(child)
        tc_pr.append(
            parse_xml(
                f'<a:{edge} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                f' w="{BORDER_EMU}" cap="flat" cmpd="sng" algn="ctr">'
                f'<a:solidFill><a:srgbClr val="{BORDER_WHITE}"/></a:solidFill>'
                f'<a:prstDash val="solid"/>'
                f'</a:{edge}>'
            )
        )


def _add_title(slide: Any, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.45)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=18, bold=True, color=NAVY)


def _add_notes(
    slide: Any, text: str, top_in: float = 6.55, height_in: float = 0.8
) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.35), Inches(top_in), Inches(12.6), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=10, color=RGBColor(0x33, 0x33, 0x33))


def _fill_cell(
    cell: Any,
    text: str,
    *,
    header: bool,
    alt: bool,
    align: Any,
) -> None:
    cell.text = ''
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    if header:
        _set_run_font(run, size_pt=TABLE_FONT_PT, bold=True, color=HEADER_FG)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
    else:
        color = SAME_FG if text == 'same' else BODY_FG
        _set_run_font(run, size_pt=TABLE_FONT_PT, bold=text == 'same', color=color)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ALT_ROW if alt else WHITE
    cell.margin_left = Emu(91440)
    cell.margin_right = Emu(91440)
    cell.margin_top = Emu(45720)
    cell.margin_bottom = Emu(45720)
    cell.text_frame.word_wrap = True
    _set_cell_borders(cell)


def _col_align(col: int, n_cols: int) -> Any:
    if n_cols >= 6:
        if col <= 1:
            return PP_ALIGN.LEFT
        return PP_ALIGN.CENTER
    if col == 0:
        return PP_ALIGN.LEFT
    return PP_ALIGN.RIGHT


def _set_column_widths(table: Any, n_cols: int, width: float) -> None:
    if n_cols >= 6:
        shares = (1.05, 1.40, 0.95, 1.05, 0.85, 0.85)
        scale = width / sum(shares)
        for j, share in enumerate(shares):
            table.columns[j].width = Inches(share * scale)
        return
    first = width * 0.34
    rest = (width - first) / max(n_cols - 1, 1)
    for j in range(n_cols):
        table.columns[j].width = Inches(first if j == 0 else rest)


def _add_table(
    slide: Any,
    grid: TableGrid,
    *,
    left: float,
    top: float,
    width: float,
) -> None:
    n_cols = len(grid.headers)
    banner = n_cols >= 6
    if not banner:
        caption = slide.shapes.add_textbox(
            Inches(left), Inches(top - 0.32), Inches(width), Inches(0.3)
        )
        tf = caption.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = grid.title
        _set_run_font(run, size_pt=TABLE_FONT_PT, bold=True, color=NAVY)

    n_header = 2 if banner else 1
    n_rows = n_header + len(grid.rows)
    height = n_header * HEADER_ROW_IN + len(grid.rows) * DATA_ROW_IN
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table
    _set_column_widths(table, n_cols, width)
    for i, row in enumerate(table.rows):
        row.height = Inches(HEADER_ROW_IN if i < n_header else DATA_ROW_IN)

    header_row = 1 if banner else 0
    if banner:
        for j in range(n_cols):
            _fill_cell(
                table.cell(0, j), '', header=True, alt=False, align=PP_ALIGN.CENTER
            )
        table.cell(0, 2).merge(table.cell(0, n_cols - 1))
        _fill_cell(
            table.cell(0, 2),
            grid.title,
            header=True,
            alt=False,
            align=PP_ALIGN.CENTER,
        )
    for j, header in enumerate(grid.headers):
        _fill_cell(
            table.cell(header_row, j),
            header,
            header=True,
            alt=False,
            align=PP_ALIGN.CENTER if n_cols >= 6 else _col_align(j, n_cols),
        )
    data_start = header_row + 1
    for i, row in enumerate(grid.rows):
        alt = i % 2 == 0
        for j, value in enumerate(row):
            _fill_cell(
                table.cell(data_start + i, j),
                value,
                header=False,
                alt=alt,
                align=_col_align(j, n_cols),
            )
    if banner and len(grid.rows) > 1:
        origin = table.cell(data_start, 0)
        origin.merge(table.cell(n_rows - 1, 0))
        origin.vertical_anchor = MSO_ANCHOR.MIDDLE


def _two_tables_slide(
    prs: Any,
    title: str,
    note: str,
    left_grid: TableGrid,
    right_grid: TableGrid,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)
    _add_notes(slide, note, top_in=0.58, height_in=2.2)
    _add_table(
        slide, left_grid, left=TABLE_LEFTS[0], top=TABLE_TOP_IN, width=TABLE_WIDTH_IN
    )
    _add_table(
        slide, right_grid, left=TABLE_LEFTS[1], top=TABLE_TOP_IN, width=TABLE_WIDTH_IN
    )


def _image_slide(prs: Any, title: str, caption: str, png: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title)
    slide.shapes.add_picture(
        str(png), Inches(0.25), Inches(0.6), width=Inches(12.8), height=Inches(6.0)
    )
    _add_notes(slide, caption, top_in=6.65)


def write_pptx(
    pair: Pair,
    top: ImplBundle,
    bottom: ImplBundle,
    out_path: Path,
    *,
    png_dir: Path,
) -> Path:
    fill_mixed_c_col(top)
    fill_mixed_c_col(bottom)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    png_dir.mkdir(parents=True, exist_ok=True)
    d_png = write_hist_png(pair, top, bottom, 'D', png_dir / f'{pair.key}_D.png')
    n_png = write_hist_png(pair, top, bottom, 'N', png_dir / f'{pair.key}_N.png')

    use_qx = pair.table_steps[-1] == 'reaggregation'
    if use_qx:
        slide1_top, slide1_bottom = qx_grids(pair, top, bottom)
        slide1_title = 'Overview of results: 221100 q and industry x'
    else:
        slide1_top, slide1_bottom = class_mwh_grids(pair, top, bottom)
        slide1_title = 'Overview of results: Electricity MWh by use class'
    d_top, d_bottom = ef_grids(pair, top, bottom, 'D')
    n_top, n_bottom = ef_grids(pair, top, bottom, 'N')

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _two_tables_slide(
        prs,
        slide1_title,
        pair.slide1_note,
        slide1_top,
        slide1_bottom,
    )
    _two_tables_slide(
        prs,
        'Overview of results: D for electricity sectors',
        pair.slide_ef_note
        + ' Electricity D EFs by disaggregation step (kg CO₂e / USD).',
        d_top,
        d_bottom,
    )
    _image_slide(
        prs,
        'Overview of results: D for electricity sectors',
        pair.slide3_caption,
        d_png,
    )
    slide4_note = pair.slide_ef_note
    if pair.slide4_extra_note:
        slide4_note = f'{slide4_note} {pair.slide4_extra_note}'
    _two_tables_slide(
        prs,
        'Overview of results: N for electricity sectors',
        slide4_note + ' Electricity N EFs by disaggregation step (kg CO₂e / USD).',
        n_top,
        n_bottom,
    )
    _image_slide(
        prs,
        'Overview of results: N for electricity sectors',
        pair.slide5_caption,
        n_png,
    )
    prs.save(str(out_path))
    return out_path
