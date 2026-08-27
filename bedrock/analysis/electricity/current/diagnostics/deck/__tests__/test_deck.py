"""Synthetic tests for the five-slide electricity comparison deck."""

from __future__ import annotations

import zipfile
from pathlib import Path
from types import SimpleNamespace

import pandas as pd
import pytest
from PIL import Image
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from bedrock.analysis.electricity.current.diagnostics.deck import histograms as hist
from bedrock.analysis.electricity.current.diagnostics.deck.data import (
    NA,
    SAME,
    ImplBundle,
    StepSnapshot,
    c_col_is_monetary,
    ef_kg_per_usd,
    fill_mixed_c_col,
    format_twh,
    format_usd,
    sector_ef_usd,
    star_aggregate,
    values_match,
)
from bedrock.analysis.electricity.current.diagnostics.deck.histograms import (
    _panel_data,
    frozen_panel_png,
    pairwise_frame,
    perc_frame,
    stack_panel_pngs,
    write_hist_png,
)
from bedrock.analysis.electricity.current.diagnostics.deck.pairs import (
    AGGREGATE_ONLY_NA,
    IMPLEMENTATIONS,
    PAIRS,
    STEPS,
    ImplId,
    config_for_step,
    na_sectors_at_step,
)
from bedrock.analysis.electricity.current.diagnostics.deck.pptx_write import write_pptx
from bedrock.analysis.electricity.current.diagnostics.deck.tables import (
    class_mwh_grids,
    ef_grids,
    format_cell_pair,
    qx_grids,
)
from bedrock.analysis.electricity.historical.original_vs_eia_anchored_deck.published import (
    published_ef,
)


def _step(
    d: dict[str, float],
    n: dict[str, float] | None = None,
    *,
    mixed: bool = False,
    c_col: float | None = None,
    class_mwh: dict[str, float] | None = None,
    target: dict[str, float] | None = None,
    q: dict[str, float] | None = None,
    x: dict[str, float] | None = None,
) -> StepSnapshot:
    return StepSnapshot(
        d=pd.Series(d, dtype=float),
        n=pd.Series(n if n is not None else d, dtype=float),
        q=pd.Series(q, dtype=float) if q is not None else None,
        x=pd.Series(x, dtype=float) if x is not None else None,
        mixed=mixed,
        c_col=c_col,
        class_mwh=pd.Series(class_mwh, dtype=float) if class_mwh is not None else None,
        class_mwh_target=pd.Series(target, dtype=float) if target is not None else None,
    )


def _eia_class() -> tuple[dict[str, float], dict[str, float]]:
    model = {
        'Residential': 1.5488e9,
        'Commercial': 1.0e9,
        'Industrial': 1.2e9,
        'Transportation': 0.1e9,
        'Exports': 0.4638e9,
    }
    return model, dict(model)


def test_ef_kg_per_usd_scales_generation_when_mixed() -> None:
    raw = pd.Series({'221110': 100.0, '221121': 0.2})
    usd = ef_kg_per_usd(raw, mixed=True, c_col=0.05)
    assert usd.loc['221110'] == pytest.approx(5.0)
    assert usd.loc['221121'] == pytest.approx(0.2)


def test_identity_c_col_does_not_convert_generation() -> None:
    raw = pd.Series({'221110': 398.866})
    usd = ef_kg_per_usd(raw, mixed=True, c_col=1.0)
    assert usd.loc['221110'] == pytest.approx(398.866)
    assert not c_col_is_monetary(1.0)
    assert c_col_is_monetary(0.0216)


def test_fill_mixed_c_col_from_three_way_dollar_q() -> None:
    q_mwh = 4.312557e9
    q_usd = 1.9951353935173175e11
    mixed = _step(
        {'221110': 398.86642559883023, '221121': 0.36, '221122': 0.119},
        n={'221110': 398.86642559883023, '221121': 0.36, '221122': 0.119},
        mixed=True,
        c_col=None,
        q={'221110': q_mwh, '221121': 1e9, '221122': 1e10},
    )
    three = _step(
        {'221110': 8.622, '221121': 0.36, '221122': 0.119},
        q={'221110': q_usd, '221121': 1e9, '221122': 1e10},
    )
    bundle = ImplBundle('eia_gtd', {'three_way': three, 'mixed_units': mixed})
    fill_mixed_c_col(bundle)
    assert mixed.c_col == pytest.approx(q_mwh / q_usd)
    assert sector_ef_usd(mixed, 'N', '221110') == pytest.approx(8.622, rel=1e-3)


def test_fill_mixed_c_col_replaces_identity() -> None:
    q_mwh = 4.312557e9
    q_usd = 1.9951353935173175e11
    mixed = _step(
        {'221110': 398.866},
        mixed=True,
        c_col=1.0,
        q={'221110': q_mwh},
    )
    three = _step({'221110': 8.622}, q={'221110': q_usd})
    fill_mixed_c_col(ImplBundle('eia_gtd', {'three_way': three, 'mixed_units': mixed}))
    assert mixed.c_col == pytest.approx(q_mwh / q_usd)


def test_fill_mixed_c_col_keeps_original_freeze_c_col() -> None:
    mixed = _step(
        {'221110': 816.64},
        mixed=True,
        c_col=0.013975279012407282,
        q={'221110': 4.312557e9},
    )
    three = _step({'221110': 10.176}, q={'221110': 3.08e11})
    fill_mixed_c_col(ImplBundle('original', {'three_way': three, 'mixed_units': mixed}))
    assert mixed.c_col == pytest.approx(0.013975279012407282)


def test_star_aggregate_x_weighted() -> None:
    step = _step(
        {'221110': 10.0, '221121': 1.0, '221122': 0.0},
        x={'221110': 2.0, '221121': 2.0, '221122': 6.0},
    )
    # (10*2 + 1*2 + 0*6) / 10 = 2.2
    assert star_aggregate(step, 'D') == pytest.approx(2.2)


def test_format_twh() -> None:
    assert format_twh(1_548_800_000) == '1,548.8 TWh'


def test_na_and_same_cells() -> None:
    left = ImplBundle(
        'mecs_mixed_units',
        {
            'footing': _step({'221100': 2.409}),
            'reallocation': _step({'221100': 2.416}),
            'three_way': _step({'221110': 8.0, '221121': 0.226, '221122': 0.0}),
            'mixed_units': _step(
                {'221110': 8.0, '221121': 0.226, '221122': 0.0},
                mixed=True,
                c_col=1.0,
            ),
        },
    )
    right = ImplBundle(
        'eia_gtd',
        {
            'footing': _step({'221100': 2.409}),
            'reallocation': _step({'221100': 2.416}),
            'three_way': _step({'221110': 8.0, '221121': 0.226, '221122': 0.0}),
            'mixed_units': _step(
                {'221110': 8.0, '221121': 0.226, '221122': 0.0},
                mixed=True,
                c_col=1.0,
            ),
        },
    )
    assert format_cell_pair(left, right, 'D', '221110', 'footing', False) == NA
    assert format_cell_pair(left, right, 'D', '221100', 'footing', True) == SAME
    assert format_cell_pair(right, left, 'D', '221110', 'three_way', True) != SAME
    assert values_match(2.409, 2.409)


def test_class_mwh_same_when_totals_match() -> None:
    model, target = _eia_class()
    mixed = _step({'221100': 1.0}, class_mwh=model, target=target)
    top = ImplBundle('mecs_mixed_units', {'mixed_units': mixed})
    bottom = ImplBundle('eia_gtd', {'mixed_units': mixed})
    pair = PAIRS['mecs_mixed_units_vs_eia_gtd']
    top_grid, bottom_grid = class_mwh_grids(pair, top, bottom)
    assert bottom_grid.rows[0][1] == SAME
    assert top_grid.rows[-1][1] == '4,312.6 TWh'


def test_pairwise_perc_formula() -> None:
    left = _step({'1111A0': 1.1, '221121': 0.4})
    right = _step({'1111A0': 1.0, '221121': 0.4})
    frame, drops = pairwise_frame(left, right, 'D')
    by_sector = frame.set_index('sector')['perc_diff']
    assert by_sector.loc['1111A0'] == pytest.approx(0.1)
    assert by_sector.loc['221121'] == pytest.approx(0.0)
    assert drops == []


def test_perc_frame_inner_join() -> None:
    frame = perc_frame(
        pd.Series({'a': 2.0, 'b': 1.0}),
        pd.Series({'a': 1.0, 'c': 9.0}),
    )
    assert list(frame['sector']) == ['a']
    assert frame['perc_diff'].iloc[0] == pytest.approx(1.0)


def test_ef_grids_two_tables() -> None:
    pair = PAIRS['mecs_mixed_units_vs_original']
    original = ImplBundle(
        'original',
        {
            'footing': _step({'221100': 2.386}),
            'three_way': _step({'221110': 7.138, '221121': 0.225, '221122': 0.0}),
        },
    )
    current = ImplBundle(
        'mecs_mixed_units',
        {
            'footing': _step({'221100': 2.409}),
            'three_way': _step({'221110': 7.197, '221121': 0.226, '221122': 0.0}),
        },
    )
    top, bottom = ef_grids(pair, original, current, 'D')
    assert top.headers[2] == 'v0.2'
    assert bottom.headers[2] == 'v0.3.1'
    assert top.rows[2][1] == '221110 (G)'
    n_top, _n_bottom = ef_grids(pair, original, current, 'N')
    assert n_top.rows[2][4] == '9.213'
    assert n_top.rows[2][5] == '10.070'


def test_stack_panel_pngs_preserves_pixels(tmp_path: Path) -> None:
    a = tmp_path / 'a.png'
    b = tmp_path / 'b.png'
    Image.new('RGB', (10, 4), (255, 0, 0)).save(a)
    Image.new('RGB', (10, 6), (0, 0, 255)).save(b)
    out = tmp_path / 'stacked.png'
    stack_panel_pngs([a, b], out)
    im = Image.open(out)
    assert im.size == (10, 10)


def test_write_hist_png_uses_frozen_vs_footing_panels(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    panel = tmp_path / 'panel'
    panel.mkdir()
    Image.new('RGB', (20, 8), (255, 0, 0)).save(
        panel / hist.FROZEN_PANEL_PNG[('original', 'D')]
    )
    Image.new('RGB', (20, 8), (0, 0, 255)).save(
        panel / hist.FROZEN_PANEL_PNG[('eia_gtd', 'D')]
    )
    monkeypatch.setattr(hist, 'FIGURES_DIR', panel)
    out = tmp_path / 'hist.png'
    write_hist_png(
        PAIRS['eia_gtd_vs_original'],
        ImplBundle('original', {}),
        ImplBundle('eia_gtd', {}),
        'D',
        out,
    )
    im = Image.open(out)
    assert im.size == (20, 16)


def test_published_original_n_matches_pptx() -> None:
    assert published_ef('original', 'N', '221110', 'three_way') == pytest.approx(9.213)
    assert published_ef('original', 'N', '221110', 'mixed_units') == pytest.approx(
        10.070
    )
    assert published_ef('eia_gtd', 'N', '221110', 'three_way') == pytest.approx(8.622)


def test_published_original_panel_is_on_disk() -> None:
    path = frozen_panel_png('original', 'D')
    assert path is not None
    assert path.name == 'v0.2_original_electricity_disagg_D.png'
    published: tuple[tuple[ImplId, str], ...] = (
        ('original', 'D'),
        ('original', 'N'),
        ('eia_gtd', 'D'),
        ('eia_gtd', 'N'),
    )
    for impl_id, kind in published:
        found = frozen_panel_png(impl_id, kind)
        assert found is not None, f'missing published panel for {impl_id} {kind}'


def test_write_pptx_five_slides(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.setattr(hist, 'FIGURES_DIR', tmp_path / 'no_frozen_panels')
    model, target = _eia_class()
    sectors = {f'{i:06d}': float(i) for i in range(10)}
    sectors.update({'221110': 8.0, '221121': 0.3, '221122': 0.1, '221100': 2.4})
    step_3way = _step(sectors, class_mwh=model, target=target)
    step_mixed = _step(sectors, mixed=True, c_col=0.02, class_mwh=model, target=target)
    step_re = _step({'221100': 2.4, **{f'{i:06d}': float(i) for i in range(10)}})
    step_foot = _step(
        {'221100': 2.39, **{f'{i:06d}': float(i) * 0.9 for i in range(10)}}
    )
    bundle_a = ImplBundle(
        'mecs_mixed_units',
        {
            'footing': step_foot,
            'reallocation': step_re,
            'three_way': step_3way,
            'mixed_units': step_mixed,
        },
    )
    bundle_b = ImplBundle(
        'eia_gtd',
        {
            'footing': step_foot,
            'reallocation': step_re,
            'three_way': step_3way,
            'mixed_units': step_mixed,
        },
    )
    out = tmp_path / 'deck.pptx'
    write_pptx(
        PAIRS['mecs_mixed_units_vs_eia_gtd'], bundle_a, bundle_b, out, png_dir=tmp_path
    )
    assert out.is_file()
    with zipfile.ZipFile(out) as z:
        slides = [n for n in z.namelist() if n.startswith('ppt/slides/slide')]
        media = [n for n in z.namelist() if n.startswith('ppt/media/')]
    assert len(slides) == 5
    assert len(media) >= 2
    prs = PptxPresentation(str(out))
    d_tables = [s for s in prs.slides[1].shapes if s.has_table]
    assert d_tables
    table = d_tables[0].table
    run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
    assert run.font.size is not None
    assert run.font.size.pt == 12
    assert d_tables[0].height < Inches(3.2)
    assert d_tables[0].width < Inches(6.5)


def test_production_pair_schema_and_na() -> None:
    pair = PAIRS['mecs_mixed_units_vs_production']
    assert pair.hist_baseline == 'peer'
    assert pair.top == 'mecs_mixed_units'
    assert pair.bottom == 'production'
    for step in ('footing', 'reallocation', 'three_way', 'mixed_units'):
        assert na_sectors_at_step('production', step) == AGGREGATE_ONLY_NA
    assert '221110' not in na_sectors_at_step('mecs_mixed_units', 'three_way')


def test_production_tables_and_peer_histograms() -> None:
    pair = PAIRS['mecs_mixed_units_vs_production']
    prod_step = _step({'221100': 2.5, '1111A0': 2.0})
    production = ImplBundle(
        'production',
        {
            'footing': prod_step,
            'reallocation': prod_step,
            'three_way': prod_step,
            'mixed_units': prod_step,
        },
    )
    current = ImplBundle(
        'mecs_mixed_units',
        {
            'footing': _step({'221100': 2.4, '1111A0': 1.0}),
            'reallocation': _step({'221100': 2.41, '1111A0': 1.0}),
            'three_way': _step(
                {'221110': 8.0, '221121': 0.2, '221122': 0.0, '1111A0': 1.0}
            ),
            'mixed_units': _step(
                {'221110': 8.0, '221121': 0.2, '221122': 0.0, '1111A0': 1.0},
                mixed=True,
                c_col=0.02,
            ),
        },
    )
    assert (
        format_cell_pair(production, current, 'D', '221110', 'three_way', False) == NA
    )
    assert format_cell_pair(
        production, current, 'D', '221100', 'mixed_units', False
    ) == ('2.500')
    assert (
        format_cell_pair(current, production, 'D', '221100', 'three_way', False) == NA
    )
    _top_class, bottom_class = class_mwh_grids(pair, current, production)
    assert bottom_class.rows[0][0] == '(no class MWh)'
    frame, _drops = _panel_data(
        pair, current, 'mecs_mixed_units', current, production, 'reallocation', 'D'
    )
    by_sector = frame.set_index('sector')['perc_diff']
    assert by_sector.loc['1111A0'] == pytest.approx(-0.5)


def test_write_pptx_mecs_mixed_units_vs_production(tmp_path: Path) -> None:
    sectors = {f'{i:06d}': float(i) for i in range(10)}
    sectors.update({'221110': 8.0, '221121': 0.3, '221122': 0.1, '221100': 2.4})
    prod_sectors = {f'{i:06d}': float(i) * 0.95 for i in range(10)}
    prod_sectors['221100'] = 2.5
    current = ImplBundle(
        'mecs_mixed_units',
        {
            'footing': _step(
                {'221100': 2.39, **{f'{i:06d}': float(i) * 0.9 for i in range(10)}}
            ),
            'reallocation': _step(
                {'221100': 2.4, **{f'{i:06d}': float(i) for i in range(10)}}
            ),
            'three_way': _step(sectors),
            'mixed_units': _step(sectors, mixed=True, c_col=0.02),
        },
    )
    production = ImplBundle(
        'production',
        {step: _step(prod_sectors) for step in STEPS},
    )
    out = tmp_path / 'mecs_mixed_units_vs_production.pptx'
    write_pptx(
        PAIRS['mecs_mixed_units_vs_production'],
        current,
        production,
        out,
        png_dir=tmp_path,
    )
    assert out.is_file()
    with zipfile.ZipFile(out) as z:
        slides = [n for n in z.namelist() if n.startswith('ppt/slides/slide')]
        media = [n for n in z.namelist() if n.startswith('ppt/media/')]
    assert len(slides) == 5
    assert len(media) >= 2


def test_no_legacy_pair_or_impl_keys() -> None:
    assert 'current' not in IMPLEMENTATIONS
    assert 'current_vs_eia_gtd' not in PAIRS
    assert 'current_vs_original' not in PAIRS
    assert 'current_vs_production' not in PAIRS
    assert 'mixed_units_vs_production' not in PAIRS
    assert 'reaggregation_vs_production' not in PAIRS
    for key in PAIRS:
        assert 'current' not in key


def test_reaggregated_pair_table_and_hist_steps() -> None:
    pair = PAIRS['reaggregated_vs_production']
    assert pair.table_steps == (
        'footing',
        'reallocation',
        'three_way',
        'reaggregation',
    )
    assert pair.hist_steps == ('reallocation', 'three_way', 'reaggregation')
    assert pair.top == 'production'
    assert pair.bottom == 'reaggregation'
    mixed_pair = PAIRS['mecs_mixed_units_vs_production']
    assert mixed_pair.table_steps[-1] == 'mixed_units'
    assert 'reaggregation' not in mixed_pair.table_steps
    assert na_sectors_at_step('reaggregation', 'reaggregation') == frozenset(
        {'221100*', '221110', '221121', '221122'}
    )
    assert '221100' not in na_sectors_at_step('reaggregation', 'reaggregation')
    assert '221110' not in na_sectors_at_step('reaggregation', 'three_way')


def test_reagg_and_mecs_share_first_three_step_configs() -> None:
    reagg = IMPLEMENTATIONS['reaggregation']
    mecs = IMPLEMENTATIONS['mecs_mixed_units']
    for step in ('footing', 'reallocation', 'three_way'):
        assert config_for_step(reagg, step) == config_for_step(mecs, step)
        assert na_sectors_at_step('reaggregation', step) == na_sectors_at_step(
            'mecs_mixed_units', step
        )
    assert reagg.industrial_weights == mecs.industrial_weights == 'mecs'
    assert config_for_step(reagg, 'reaggregation') != config_for_step(
        mecs, 'mixed_units'
    )


def test_qx_grids_use_reaggregation_step_x() -> None:
    pair = PAIRS['reaggregated_vs_production']
    production = ImplBundle(
        'production',
        {
            'reaggregation': _step(
                {'221100': 2.5}, q={'221100': 1.0e11}, x={'221100': 2.0e11}
            )
        },
    )
    reagg = ImplBundle(
        'reaggregation',
        {
            'reaggregation': _step(
                {'221100': 2.6}, q={'221100': 1.2e11}, x={'221100': 2.1e11}
            )
        },
    )
    top, bottom = qx_grids(pair, production, reagg)
    assert top.headers == ('Sector', 'q (USD)', 'x (USD)')
    assert top.rows[0][1] == format_usd(1.0e11)
    assert top.rows[0][2] == format_usd(2.0e11)
    assert bottom.rows[0][1] == format_usd(1.2e11)


def test_load_impl_bundle_uses_pair_table_steps(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from bedrock.analysis.electricity.current.diagnostics.deck import (  # noqa: PLC0415
        sources,
    )

    def fake_load(_impl: object, step_id: str) -> StepSnapshot:
        return _step({'221100': 1.0}, q={'221100': 1.0e11}, x={'221100': 2.0e11})

    monkeypatch.setattr(sources, 'load_step', fake_load)
    pair = PAIRS['reaggregated_vs_production']
    reagg = sources.load_impl_bundle(
        IMPLEMENTATIONS['reaggregation'], table_steps=pair.table_steps
    )
    assert set(reagg.steps) == set(pair.table_steps)
    assert 'mixed_units' not in reagg.steps
    production = sources.load_impl_bundle(
        IMPLEMENTATIONS['production'], table_steps=pair.table_steps
    )
    assert 'reaggregation' in production.steps
    assert 'mixed_units' not in production.steps
    assert production.steps['reaggregation'].x is not None
    assert production.steps['reaggregation'].q is not None
    mixed_pair = PAIRS['mecs_mixed_units_vs_production']
    production_mixed = sources.load_impl_bundle(
        IMPLEMENTATIONS['production'], table_steps=mixed_pair.table_steps
    )
    assert 'mixed_units' in production_mixed.steps
    assert 'reaggregation' not in production_mixed.steps


def test_derive_step_first_three_dn_match_mecs_and_write_x(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path
) -> None:
    """Ladder steps share YAML flags; reagg impl does not collapse until last step.

    Fake D/N are keyed off the live YAML flags so a wrong Aq branch or config
    would make the first three 221100 columns diverge from ``mecs_mixed_units``.
    """
    from bedrock.analysis.electricity.current.diagnostics.deck import (  # noqa: PLC0415
        sources,
    )
    from bedrock.utils.config.usa_config import (  # noqa: PLC0415
        get_usa_config,
        reset_usa_config,
    )

    aq_scaled_calls = {'n': 0}
    aq_reagg_calls = {'n': 0}

    def _fake_scaled() -> SimpleNamespace:
        aq_scaled_calls['n'] += 1
        return SimpleNamespace(
            scaled_q=pd.Series({'221100': 1.0e11}, dtype=float, name='q')
        )

    def _fake_reagg() -> SimpleNamespace:
        aq_reagg_calls['n'] += 1
        return SimpleNamespace(
            scaled_q=pd.Series({'221100': 1.2e11}, dtype=float, name='q')
        )

    def _fake_efs() -> SimpleNamespace:
        cfg = get_usa_config()
        if cfg.implement_electricity_reaggregation:
            d = pd.Series({'221100': 2.55, '1111A0': 1.0}, dtype=float)
        elif cfg.implement_electricity_disaggregation:
            d = pd.Series(
                {'221110': 8.0, '221121': 0.2, '221122': 0.1, '1111A0': 1.0},
                dtype=float,
            )
        elif cfg.implement_electricity_reallocation:
            d = pd.Series({'221100': 2.41, '1111A0': 1.0}, dtype=float)
        else:
            d = pd.Series({'221100': 2.40, '1111A0': 1.0}, dtype=float)
        return SimpleNamespace(
            D_new=d.rename('D').to_frame(),
            N_new=(d * 1.1).rename('N').to_frame(),
        )

    def _fake_x() -> pd.Series[float]:
        return pd.Series({'221100': 2.0e11}, dtype=float, name='x')

    monkeypatch.setattr(
        sources, 'impl_cache_dir', lambda impl, cfg: tmp_path / impl / cfg
    )
    monkeypatch.setattr(sources, '_write_class_mwh', lambda _folder: None)
    monkeypatch.setattr(
        'bedrock.transform.eeio.derived_cornerstone.derive_cornerstone_Aq_scaled',
        _fake_scaled,
    )
    monkeypatch.setattr(
        'bedrock.transform.eeio.derived_cornerstone.'
        'derive_cornerstone_Aq_reaggregated',
        _fake_reagg,
    )
    monkeypatch.setattr(
        'bedrock.utils.validation.diagnostics_helpers.pull_efs_for_diagnostics',
        _fake_efs,
    )
    monkeypatch.setattr('bedrock.publish.model_objects.get_x', _fake_x)

    reagg = IMPLEMENTATIONS['reaggregation']
    mecs = IMPLEMENTATIONS['mecs_mixed_units']
    try:
        for step in ('footing', 'reallocation', 'three_way'):
            aq_scaled_calls['n'] = 0
            aq_reagg_calls['n'] = 0
            mixed_snap = sources.derive_step(mecs, step)
            assert aq_scaled_calls['n'] == 1
            assert aq_reagg_calls['n'] == 0
            aq_scaled_calls['n'] = 0
            reagg_snap = sources.derive_step(reagg, step)
            assert aq_scaled_calls['n'] == 1
            assert aq_reagg_calls['n'] == 0
            assert mixed_snap.x is not None
            assert reagg_snap.x is not None
            assert mixed_snap.d is not None and reagg_snap.d is not None
            assert mixed_snap.n is not None and reagg_snap.n is not None
            if step == 'three_way':
                assert '221100' not in mixed_snap.d.index
                assert '221100' not in reagg_snap.d.index
            else:
                pd.testing.assert_series_equal(mixed_snap.d, reagg_snap.d)
                pd.testing.assert_series_equal(mixed_snap.n, reagg_snap.n)
                assert float(reagg_snap.d.loc['221100']) == pytest.approx(
                    2.40 if step == 'footing' else 2.41
                )

        aq_scaled_calls['n'] = 0
        aq_reagg_calls['n'] = 0
        collapsed = sources.derive_step(reagg, 'reaggregation')
        assert aq_reagg_calls['n'] == 1
        assert aq_scaled_calls['n'] == 0
        assert collapsed.x is not None
        assert collapsed.d is not None
        assert float(collapsed.d.loc['221100']) == pytest.approx(2.55)
    finally:
        reset_usa_config(should_reset_env_var=True)


def test_reagg_ef_grid_headers_and_shared_ladder() -> None:
    pair = PAIRS['reaggregated_vs_production']
    shared_foot = _step({'221100': 2.4, '1111A0': 1.0})
    shared_re = _step({'221100': 2.41, '1111A0': 1.0})
    shared_3way = _step({'221110': 8.0, '221121': 0.2, '221122': 0.0, '1111A0': 1.0})
    reagg = ImplBundle(
        'reaggregation',
        {
            'footing': shared_foot,
            'reallocation': shared_re,
            'three_way': shared_3way,
            'reaggregation': _step({'221100': 2.55, '1111A0': 1.0}),
        },
    )
    mecs = ImplBundle(
        'mecs_mixed_units',
        {
            'footing': shared_foot,
            'reallocation': shared_re,
            'three_way': shared_3way,
            'mixed_units': _step(
                {'221110': 8.0, '221121': 0.2, '221122': 0.0, '1111A0': 1.0},
                mixed=True,
                c_col=0.02,
            ),
        },
    )
    production = ImplBundle(
        'production',
        {step: _step({'221100': 2.5, '1111A0': 2.0}) for step in pair.table_steps},
    )
    top, bottom = ef_grids(pair, production, reagg, 'D')
    assert top.headers[-1] == 'reaggregation'
    assert bottom.headers[-1] == 'reaggregation'
    dummy = ImplBundle('production', {s: shared_foot for s in STEPS})
    for kind in ('D', 'N'):
        for step in ('footing', 'reallocation'):
            assert format_cell_pair(reagg, dummy, kind, '221100', step, False) == (
                format_cell_pair(mecs, dummy, kind, '221100', step, False)
            )
        assert format_cell_pair(reagg, dummy, kind, '221100', 'three_way', False) == NA
        assert format_cell_pair(mecs, dummy, kind, '221100', 'three_way', False) == NA
        assert format_cell_pair(reagg, dummy, kind, '221110', 'three_way', False) == (
            format_cell_pair(mecs, dummy, kind, '221110', 'three_way', False)
        )


def test_write_pptx_reaggregated_vs_production(tmp_path: Path) -> None:
    sectors = {f'{i:06d}': float(i) for i in range(10)}
    sectors.update({'221100': 2.55, '1111A0': 1.0})
    reagg = ImplBundle(
        'reaggregation',
        {
            'footing': _step({'221100': 2.4, **sectors}),
            'reallocation': _step({'221100': 2.41, **sectors}),
            'three_way': _step(
                {'221110': 8.0, '221121': 0.2, '221122': 0.0, **sectors}
            ),
            'reaggregation': _step(
                {'221100': 2.55, **sectors},
                q={'221100': 1.2e11},
                x={'221100': 2.1e11},
            ),
        },
    )
    production = ImplBundle(
        'production',
        {
            step: _step(
                {'221100': 2.5, **{f'{i:06d}': float(i) * 0.95 for i in range(10)}},
                q={'221100': 1.0e11},
                x={'221100': 2.0e11},
            )
            for step in PAIRS['reaggregated_vs_production'].table_steps
        },
    )
    out = tmp_path / 'reaggregated_vs_production.pptx'
    write_pptx(
        PAIRS['reaggregated_vs_production'],
        production,
        reagg,
        out,
        png_dir=tmp_path,
    )
    assert out.is_file()
    prs = PptxPresentation(str(out))
    assert len(prs.slides) == 5
    title = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text
    assert 'q and industry x' in title
