"""Assemble CEDA electricity diagnostics PPTX decks."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from bedrock.analysis.electricity.current.diagnostics.ceda_electricity.data import (
    ef_focus_row,
    load_bly_cs,
    load_config_summary,
    load_d,
    load_n,
)
from bedrock.analysis.electricity.current.diagnostics.ceda_electricity.histograms import (
    write_ladder_hist_panels,
    write_manual_adj_compare,
    write_nd_hist_pair,
)
from bedrock.analysis.electricity.current.diagnostics.ceda_electricity.paths import (
    PNG_DIR,
    ensure_ceda_dirs,
)
from bedrock.analysis.electricity.current.diagnostics.ceda_electricity.sheets import (
    DECK_A_FILENAME,
    DECK_B_FILENAME,
    ELECTRICITY_SECTOR,
    LADDER_STEPS,
    REAGG_CONFIG,
    REAGG_LABEL,
    REAGG_SHEET_ID,
)
from bedrock.analysis.electricity.current.diagnostics.ceda_electricity.waterfalls import (
    write_bly_ladder_waterfall,
    write_da_effect_scatter,
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
HEADER_BG = RGBColor(0x00, 0x5B, 0x7F)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
BODY_FG = RGBColor(0x22, 0x22, 0x22)
ALT_ROW = RGBColor(0xEB, 0xEB, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_run_font(
    run: Any, *, size_pt: float, bold: bool = False, color: RGBColor | None = None
) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = 'Calibri'
    if color is not None:
        run.font.color.rgb = color


def _add_title(slide: Any, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.35), Inches(0.12), Inches(12.6), Inches(0.45)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=18, bold=True, color=NAVY)


def _add_notes(slide: Any, text: str, *, top_in: float = 6.65) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.35), Inches(top_in), Inches(12.6), Inches(0.7)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=10, color=RGBColor(0x33, 0x33, 0x33))


def _add_bullets(slide: Any, lines: list[str], *, top_in: float = 0.7) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top_in), Inches(12.2), Inches(5.6)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _set_run_font(run, size_pt=14, color=BODY_FG)
        p.space_after = Pt(10)


def _add_picture(slide: Any, path: Path, *, top_in: float = 0.55) -> None:
    slide.shapes.add_picture(
        str(path), Inches(0.35), Inches(top_in), width=Inches(12.6)
    )


def _fmt(v: float | str | None, *, kind: str = 'ef') -> str:
    if v is None:
        return '—'
    if isinstance(v, str):
        return v
    if kind == 'pct':
        return f'{v:+.2f}%'
    if kind == 'ef':
        return f'{v:.3f}'
    return f'{v:,.2f}'


def _add_focus_table(slide: Any, rows: list[dict[str, float | str | None]]) -> None:
    headers = [
        'Country',
        'Sector',
        'N new',
        'N old infl',
        'N %',
        'N D-eff',
        'N A-eff',
        'D new',
        'D old infl',
        'D %',
    ]
    table = slide.shapes.add_table(
        len(rows) + 1,
        len(headers),
        Inches(0.35),
        Inches(0.7),
        Inches(12.6),
        Inches(0.4 + 0.32 * len(rows)),
    ).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ''
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        _set_run_font(run, size_pt=11, bold=True, color=HEADER_FG)
    keys = [
        'country',
        'sector',
        'N_new',
        'N_old_infl',
        'N_perc',
        'N_d_effect',
        'N_a_effect',
        'D_new',
        'D_old_infl',
        'D_perc',
    ]
    kinds = [
        'str',
        'str',
        'ef',
        'ef',
        'pct',
        'pct',
        'pct',
        'ef',
        'ef',
        'pct',
    ]
    for i, row in enumerate(rows):
        for j, (key, kind) in enumerate(zip(keys, kinds, strict=True)):
            cell = table.cell(i + 1, j)
            cell.text = ''
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i % 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ALT_ROW
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            val = row.get(key)
            run.text = str(val) if kind == 'str' else _fmt(val, kind=kind)  # type: ignore[arg-type]
            _set_run_font(run, size_pt=11, color=BODY_FG)


def write_deck_a(*, refresh: bool = False, out_dir: Path | None = None) -> Path:
    ensure_ceda_dirs()
    dest_dir = out_dir or ensure_ceda_dirs()
    png = PNG_DIR / 'deck_a'
    png.mkdir(parents=True, exist_ok=True)

    cfg = load_config_summary(REAGG_SHEET_ID, refresh=refresh)
    n = load_n(REAGG_SHEET_ID, refresh=refresh)
    d = load_d(REAGG_SHEET_ID, refresh=refresh)

    global_nd = write_nd_hist_pair(
        n, d, png / 'global_nd.png', scope='global', title_prefix='Global'
    )
    usa_nd = write_nd_hist_pair(
        n, d, png / 'usa_nd.png', scope='USA', title_prefix='USA'
    )
    ma_cmp = write_manual_adj_compare(n, png / 'usa_ma_compare.png', scope='USA')
    scatter = write_da_effect_scatter(n, png / 'usa_da_scatter.png', scope='USA')

    focus_countries = ('USA', 'CHN', 'GBR', 'DEU', 'JPN')
    focus_rows = [
        ef_focus_row(n, d, country=c, sector=ELECTRICITY_SECTOR)
        for c in focus_countries
    ]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(blank)
    _add_title(s0, 'CEDA Deck A: electricity reaggregation vs CEDA v8.1')
    _add_bullets(
        s0,
        [
            f'Run: {REAGG_LABEL}',
            f'Config: {cfg.get("config_name", REAGG_CONFIG)}',
            f'Sheet: {REAGG_SHEET_ID}',
            f'ceda_base_year: {cfg.get("ceda_base_year", "—")}',
            'Baseline in each sheet: CEDA v8.1 no-manual-adj (inflated) and with-manual-adj.',
            'Focus sector: 221100 Electric power generation, transmission, and distribution.',
        ],
    )

    s1 = prs.slides.add_slide(blank)
    _add_title(s1, 'Electricity sector (221100): N and D vs CEDA v8.1')
    _add_focus_table(s1, focus_rows)
    _add_notes(
        s1,
        'N/D % and D/A effects are vs no-manual-adj inflated baseline. '
        'D-effect / A-effect partition the N % change.',
    )

    s2 = prs.slides.add_slide(blank)
    _add_title(s2, 'Global N and D % diff vs CEDA v8.1 (no manual adj)')
    _add_picture(s2, global_nd)
    _add_notes(s2, 'All country×sector rows. Same histogram style as bedrock EF panels.')

    s3 = prs.slides.add_slide(blank)
    _add_title(s3, 'USA N and D % diff vs CEDA v8.1 (no manual adj)')
    _add_picture(s3, usa_nd)

    s4 = prs.slides.add_slide(blank)
    _add_title(s4, 'USA N: no-manual-adj vs with-manual-adj baselines')
    _add_picture(s4, ma_cmp)

    s5 = prs.slides.add_slide(blank)
    _add_title(s5, 'USA N change decomposition: D-effect vs A-effect')
    _add_picture(s5, scatter, top_in=0.55)
    _add_notes(
        s5,
        'Each point is one USA sector. X = direct-intensity contribution to ΔN; '
        'Y = structure (A/L) contribution.',
    )

    out = dest_dir / DECK_A_FILENAME
    prs.save(out)
    return out


def write_deck_b(*, refresh: bool = False, out_dir: Path | None = None) -> Path:
    ensure_ceda_dirs()
    dest_dir = out_dir or ensure_ceda_dirs()
    png = PNG_DIR / 'deck_b'
    png.mkdir(parents=True, exist_ok=True)

    n_frames: dict[str, Any] = {}
    d_frames: dict[str, Any] = {}
    bly_frames: dict[str, Any] = {}
    configs: dict[str, dict[str, str]] = {}
    for step in LADDER_STEPS:
        n_frames[step.key] = load_n(step.sheet_id, refresh=refresh)
        d_frames[step.key] = load_d(step.sheet_id, refresh=refresh)
        bly_frames[step.key] = load_bly_cs(step.sheet_id, refresh=refresh)
        configs[step.key] = load_config_summary(step.sheet_id, refresh=refresh)

    global_n = write_ladder_hist_panels(
        LADDER_STEPS, n_frames, png / 'global_n_panels.png', kind='N', scope='global'
    )
    usa_n = write_ladder_hist_panels(
        LADDER_STEPS, n_frames, png / 'usa_n_panels.png', kind='N', scope='USA'
    )
    global_d = write_ladder_hist_panels(
        LADDER_STEPS, d_frames, png / 'global_d_panels.png', kind='D', scope='global'
    )
    usa_d = write_ladder_hist_panels(
        LADDER_STEPS, d_frames, png / 'usa_d_panels.png', kind='D', scope='USA'
    )
    global_wf = write_bly_ladder_waterfall(
        LADDER_STEPS, bly_frames, png / 'global_bly_waterfall.png', scope='global'
    )
    usa_wf = write_bly_ladder_waterfall(
        LADDER_STEPS, bly_frames, png / 'usa_bly_waterfall.png', scope='USA'
    )

    # USA 221100 callout: g5 ABSR vs g5+elec vs Deck A elec-only.
    reagg_n = load_n(REAGG_SHEET_ID, refresh=refresh)
    reagg_d = load_d(REAGG_SHEET_ID, refresh=refresh)
    g5 = next(s for s in LADDER_STEPS if s.key == 'g5')
    g5e = next(s for s in LADDER_STEPS if s.key == 'g5e')
    callout_rows = [
        {
            **ef_focus_row(
                n_frames[g5.key],
                d_frames[g5.key],
                country='USA',
                sector=ELECTRICITY_SECTOR,
            ),
            'country': 'USA g5 ABSR',
        },
        {
            **ef_focus_row(
                n_frames[g5e.key],
                d_frames[g5e.key],
                country='USA',
                sector=ELECTRICITY_SECTOR,
            ),
            'country': 'USA g5+elec',
        },
        {
            **ef_focus_row(reagg_n, reagg_d, country='USA', sector=ELECTRICITY_SECTOR),
            'country': 'USA elec-only',
        },
    ]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(blank)
    _add_title(s0, 'CEDA Deck B: cumulative g1→g5 ABSR + electricity')
    lines = [
        'Each rung adds one change on the previous state (CEDA v0.4 cumulative ladder).',
        'g5e = g5 ABSR + electricity improvements (reads as the electricity effect on the ladder).',
        '',
    ]
    for step in LADDER_STEPS:
        cfg_name = configs[step.key].get('config_name', step.title)
        lines.append(f'{step.key}: {cfg_name}')
    _add_bullets(s0, lines)

    s1 = prs.slides.add_slide(blank)
    _add_title(s1, 'Global N % diff vs CEDA v8.1 — one panel per ladder step')
    _add_picture(s1, global_n, top_in=0.55)
    _add_notes(
        s1,
        'In-sheet N_no_manual_adj_perc_diff at each cumulative state. '
        'Compare g5 vs g5e for the electricity-improvements increment.',
    )

    s2 = prs.slides.add_slide(blank)
    _add_title(s2, 'USA N % diff vs CEDA v8.1 — one panel per ladder step')
    _add_picture(s2, usa_n, top_in=0.55)
    _add_notes(s2, 'Last panel (g5e) adds electricity on top of g5 ABSR.')

    s3 = prs.slides.add_slide(blank)
    _add_title(s3, 'Global D % diff vs CEDA v8.1 — one panel per ladder step')
    _add_picture(s3, global_d, top_in=0.55)

    s4 = prs.slides.add_slide(blank)
    _add_title(s4, 'USA D % diff vs CEDA v8.1 — one panel per ladder step')
    _add_picture(s4, usa_d, top_in=0.55)

    s5 = prs.slides.add_slide(blank)
    _add_title(s5, 'Global BLy net waterfall (MtCO2e), g1→g5e')
    _add_picture(s5, global_wf, top_in=0.55)
    _add_notes(
        s5,
        'Start bar = g1 total BLy; green/red = net step deltas; end bar = g5e total. '
        'The Δg5e bar is the electricity-improvements effect after ABSR.',
    )

    s6 = prs.slides.add_slide(blank)
    _add_title(s6, 'USA BLy net waterfall (MtCO2e), g1→g5e')
    _add_picture(s6, usa_wf, top_in=0.55)
    _add_notes(s6, 'Δg5e = USA BLy change from adding electricity improvements on g5.')

    s7 = prs.slides.add_slide(blank)
    _add_title(s7, 'USA 221100 callout: g5 vs g5+elec vs elec-only (Deck A)')
    _add_focus_table(s7, callout_rows)
    _add_notes(
        s7,
        'g5 = ABSR only. g5+elec = ABSR + electricity improvements. '
        'elec-only = Deck A (electricity on v0.3 settings, no ABSR ladder).',
    )

    out = dest_dir / DECK_B_FILENAME
    prs.save(out)
    return out
